import csv
import six
import copy
import pysftp
import shutil
import os.path
import requests
from PIL import Image
from pptx import Presentation


class CsvToPptx:
    """
    Main class for modify PPTX-file with used CSV-file(input data).
    And send result - SFTP.
    """

    dir_base = os.path.dirname(os.path.abspath(__file__))
    prs = None
    CsvFilenameInput: str = ""
    PptxFilenameInput: str = ""
    PptxFilenameOutput: str = ""
    sftp_host: str = "127.0.0.1"
    sftp_port: int = 22
    sftp_username: str = "user"
    sftp_password: str = "password"
    sftp_dir_target: str = "/public"

    def __init__(self,
                 csv_filename_input: str = "input.csv",
                 pptx_filename_input: str = "input.pptx",
                 pptx_filename_output: str = "output.pptx"):
        """
        :param PptxFilenameInput: this file is readable
        :param PptxFilenameOutput: this file is writeable
        """
        self.set_csv_filename_input(csv_filename_input)
        self.set_pptx_filename_input(pptx_filename_input)
        self.set_pptx_filename_output(pptx_filename_output)

    # PPTX
    def set_pptx_filename_input(self, filename: str = "input.pptx") -> None:
        """
        :param filename: use another file to read
        :return: None
        """
        self.PptxFilenameInput = filename
        self.prs = Presentation(self.PptxFilenameInput)

    def set_pptx_filename_output(self, filename: str = "output.pptx") -> None:
        """
        :param filename: use another file to write
        :return: None
        """
        self.PptxFilenameOutput = filename

    def save_pptx(self) -> None:
        """
        Save file with current name <<PptxFilenameOutput>>
        :return:
        """
        self.prs.save(self.PptxFilenameOutput)

    def pptx_extract_all_text_from_sliders(self) -> list():
        """
        Method name.
        :return: list of lists of text
        """
        text_runs = []

        for slide in self.prs.slides:
            print(slide)
            for shape in slide.shapes:
                print("    ", shape)
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
                        print("        ", run.text)
        return text_runs

    def pptx_delete_placeholder_in_slide(self,
                                         id_slide: int = 0,
                                         id_shape: int = 0) -> None:
        """
        Replace image
        :param id_slide: Slide number(ID)
        :param id_shape: Shape number(ID) into Slide
        :return: None
        """
        image_into_pptx = self.prs.slides[id_slide].shapes[id_shape]
        pic_into_pptx = image_into_pptx._element
        pic_into_pptx.getparent().remove(pic_into_pptx)

    def pptx_replace_text_in_slide(self,
                                   id_slide: int = 0,
                                   id_shape: int = 0,
                                   new_text: str = "new"):
        text_into_pptx = self.prs.slides[id_slide].shapes[id_shape]
        if not text_into_pptx.has_text_frame:
            return

        for paragraph in text_into_pptx.text_frame.paragraphs:
            for run in paragraph.runs:
                run.text = new_text

    def pptx_replace_image_in_slide(self,
                                    id_slide: int = 0,
                                    id_shape: int = 0,
                                    filename_img: str = "new_image.png",
                                    filename_default: str = "input/default.png") -> None:
        """
        Replace image
        :param id_slide: Slide number(ID)
        :param id_shape: Shape number(ID) into Slide
        :param filename_img: Path to new image-file.
        :param filename_default: DEFAULT: Path to new image-file if filename_img not exist.
        :return: None
        """
        if os.path.isfile(filename_default) == False:
            print("You must have a DEFAULT-file - 'input/default.png'")
            return
        if os.path.isfile(filename_img) == False:
            filename_img = filename_default
        with Image.open(filename_img) as im:
            w, h = im.size
        picture_old = self.prs.slides[id_slide].shapes[id_shape]
        picture_new = self.prs.slides[id_slide].shapes.add_picture(
            filename_img,
            left=picture_old.left,
            top=(picture_old.top + picture_old.height / 2)
                - (picture_old.width * h / w) / 2,
            width=picture_old.width)
        pic_old = picture_old._element
        pic_new = picture_new._element
        pic_new.parent = pic_old.getparent()
        pic_old.addnext(pic_new)
        pic_old.getparent().remove(pic_old)

    def pptx_get_image_from_url(self, image_url: str = "http"):
        try:
            os.stat(self.dir_base+"/output/")
        except:
            os.mkdir(self.dir_base+"/output/")
        filename = image_url.split("/")[-1]
        r = requests.get(image_url, stream = True)# Check if the image was retrieved successfully
        if r.status_code == 200:
            r.raw.decode_content = True
            with open(self.dir_base+"/output/" + filename, 'wb') as f:
                shutil.copyfileobj(r.raw, f)
        else:
            pass
        return self.dir_base+"/output/" + filename

    # CSV
    def set_csv_filename_input(self, filename: str = "input.csv") -> None:
        """
        :param filename: use another file to read
        :return: None
        """
        self.CsvFilenameInput = filename

    def csv_extract_all_text_from_table(self) -> list():
        """
        Method name.
        :return: list of lists of text
        """
        text_rows = []

        with open(self.CsvFilenameInput, 'r', encoding='UTF-8', newline='') as csvfile:
            spamreader = csv.reader(csvfile, delimiter=',', quotechar='|')
            for row in spamreader:
                text_rows.append(row)
        return text_rows

    def csv_extract_row(self, number_row: int = 0) -> list():
        """
        :param number_row: number ROW
        :return: list of lists of text
        """
        with open(self.CsvFilenameInput, 'r', encoding='UTF-8', newline='') as csvfile:
            spamreader = csv.reader(csvfile, delimiter=',', quotechar='|')
            for idx, row in enumerate(spamreader):
                if idx == number_row:
                    return row
        return []

    # SFTP
    def set_sftp_host(self, set_host: str = "127.0.0.1") -> None:
        self.sftp_host = set_host

    def set_sftp_port(self, set_port: int = 22) -> None:
        self.sftp_port = set_port

    def set_sftp_host_port(self,
                           set_host: str = "127.0.0.1",
                           set_port: int = 22) -> None:
        self.set_sftp_host(set_host=set_host)
        self.set_sftp_port(set_port=set_port)

    def set_sftp_username(self, set_username: str = "user") -> None:
        self.sftp_username = set_username

    def set_sftp_password(self, set_password: str = "qwerty") -> None:
        self.sftp_password = set_password

    def set_sftp_username_password(self,
                                   set_username: str = "user",
                                   set_password: str = "qwerty") -> None:
        self.set_sftp_username(set_username=set_username)
        self.set_sftp_password(set_password=set_password)

    def set_sftp_dir_target(self, set_dir_target: str = "/public") -> None:
        self.sftp_dir_target = set_dir_target

    def send_file_to_sftp(self,
                          set_host: str = "127.0.0.1",
                          set_port: int = 22,
                          set_username: str = "user",
                          set_password: str = "user",
                          set_dir_target: str = "/public",
                          set_file_for_send: str = "input.txt") -> None:
        self.set_sftp_host_port(set_host=set_host, set_port=set_port)
        self.set_sftp_username_password(set_username=set_username, set_password=set_password)
        self.set_sftp_dir_target(set_dir_target=set_dir_target)
        cnopts = pysftp.CnOpts(knownhosts=os.getenv("HOME")+'/.ssh/known_hosts')
        cnopts.hostkeys = None
        with pysftp.Connection(self.sftp_host, port=self.sftp_port,
                               username=self.sftp_username,
                               password=self.sftp_password,
                               cnopts=cnopts) as sftp:
            with sftp.cd(self.sftp_dir_target):
                sftp.put(set_file_for_send)

    def send_file_to(self, set_file_for_send: str = "input.txt") -> None:
        self.send_file_to_sftp(set_host=self.sftp_host,
                               set_port=self.sftp_port,
                               set_username=self.sftp_username,
                               set_password=self.sftp_password,
                               set_dir_target=self.sftp_dir_target,
                               set_file_for_send=set_file_for_send)

    # AUTOMATE
